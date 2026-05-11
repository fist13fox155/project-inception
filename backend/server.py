"""
Project Inception - Civilian AI Assistant Backend
- JARVIS Chat (Claude Sonnet 4.5 via Emergent Universal Key)
- Document Generation (PDF/PPTX from AI prompts)
- Stock Tracker (Alpha Vantage real-time + AI buy/sell recommendations)
- TTS / STT via OpenAI through Emergent gateway
"""
import os
import io
import re
import json
import uuid
import base64
import asyncio
import logging
from pathlib import Path
from datetime import datetime, timezone
from typing import List, Optional, Literal

import httpx
from fastapi import FastAPI, APIRouter, HTTPException, UploadFile, File, Form, WebSocket, WebSocketDisconnect
from fastapi.responses import Response, FileResponse
from dotenv import load_dotenv
from starlette.middleware.cors import CORSMiddleware
from motor.motor_asyncio import AsyncIOMotorClient
from pydantic import BaseModel, Field

from emergentintegrations.llm.chat import LlmChat, UserMessage

# ReportLab for PDF
from reportlab.lib.pagesizes import LETTER
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib.colors import HexColor
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib.enums import TA_LEFT, TA_CENTER

# python-pptx for PowerPoint
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

ROOT_DIR = Path(__file__).parent
load_dotenv(ROOT_DIR / ".env")

# ---------- Mongo ----------
mongo_url = os.environ["MONGO_URL"]
mongo_client = AsyncIOMotorClient(mongo_url)
db = mongo_client[os.environ["DB_NAME"]]

# ---------- Environment ----------
EMERGENT_LLM_KEY = os.environ.get("EMERGENT_LLM_KEY", "")
ALPHA_VANTAGE_KEY = os.environ.get("ALPHA_VANTAGE_KEY", "")
FINNHUB_KEY = os.environ.get("FINNHUB_KEY", "")
CLAUDE_MODEL = "claude-sonnet-4-5-20250929"

# OpenAI client for TTS / Whisper via Emergent gateway
from openai import AsyncOpenAI
oai_client = AsyncOpenAI(
    api_key=EMERGENT_LLM_KEY,
    base_url="https://integrations.emergentagent.com/llm",
)

logger = logging.getLogger("inception")
logging.basicConfig(level=logging.INFO, format="%(asctime)s %(levelname)s %(message)s")

app = FastAPI(title="Project Inception API")
api = APIRouter(prefix="/api")

# ==================================================================
# JARVIS Chat
# ==================================================================
JARVIS_SYSTEM = (
    "You are JARVIS — the personal AI assistant for the user (the 'Architect') in the "
    "civilian application 'Project Inception'. Speak with calm precision, a touch of British wit, "
    "and Iron-Man inspired formality. Be concise (2-4 sentences unless asked for depth). "
    "You help with: tracking stocks, market analysis, when to buy/sell with reasoning, "
    "world news that affects tracked stocks, generating PDF and PowerPoint documents from ideas, "
    "and everyday errands. Never give blind financial guarantees — always note risk. "
    "Address the user as 'Architect' or 'sir/ma'am' occasionally."
)


class ChatRequest(BaseModel):
    session_id: str
    message: str


class ChatMessage(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    session_id: str
    role: Literal["user", "assistant"]
    content: str
    timestamp: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))


@api.post("/chat")
async def chat_jarvis(req: ChatRequest):
    if not EMERGENT_LLM_KEY:
        raise HTTPException(500, "EMERGENT_LLM_KEY missing")

    # Save user msg
    user_msg = ChatMessage(session_id=req.session_id, role="user", content=req.message)
    await db.chat_messages.insert_one(user_msg.model_dump())

    chat = LlmChat(
        api_key=EMERGENT_LLM_KEY,
        session_id=req.session_id,
        system_message=JARVIS_SYSTEM,
    ).with_model("anthropic", CLAUDE_MODEL)

    try:
        reply = await chat.send_message(UserMessage(text=req.message))
    except Exception as e:
        logger.exception("chat failed")
        raise HTTPException(500, f"AI error: {e}")

    asst_msg = ChatMessage(session_id=req.session_id, role="assistant", content=reply)
    await db.chat_messages.insert_one(asst_msg.model_dump())

    return {"reply": reply, "message_id": asst_msg.id, "timestamp": asst_msg.timestamp.isoformat()}


@api.get("/chat/history/{session_id}")
async def chat_history(session_id: str):
    docs = await db.chat_messages.find(
        {"session_id": session_id}, {"_id": 0}
    ).sort("timestamp", 1).to_list(500)
    return {"messages": docs}


# ==================================================================
# Stocks (Alpha Vantage)
# ==================================================================
ALPHA_BASE = "https://www.alphavantage.co/query"
DEFAULT_WATCHLIST = ["AAPL", "TSLA", "NVDA", "MSFT", "GOOGL"]

# In-memory short cache to respect free-tier rate limits (25/day, 5/min)
_av_cache: dict = {}
_av_cache_ttl = 60  # seconds


async def _av_get(params: dict) -> dict:
    key = json.dumps(params, sort_keys=True)
    now = datetime.now(timezone.utc).timestamp()
    cached = _av_cache.get(key)
    if cached and (now - cached["t"]) < _av_cache_ttl:
        return cached["data"]
    params["apikey"] = ALPHA_VANTAGE_KEY
    async with httpx.AsyncClient(timeout=20) as c:
        r = await c.get(ALPHA_BASE, params=params)
        r.raise_for_status()
        data = r.json()
    _av_cache[key] = {"t": now, "data": data}
    return data


def _mock_quote(symbol: str) -> dict:
    """Fallback when Alpha Vantage rate-limited."""
    import hashlib, random
    seed = int(hashlib.md5(symbol.encode()).hexdigest()[:8], 16)
    rng = random.Random(seed + int(datetime.now(timezone.utc).timestamp() // 60))
    base = 50 + (seed % 400)
    price = round(base + rng.uniform(-5, 5), 2)
    change_pct = round(rng.uniform(-3, 3), 2)
    spark = [round(price + rng.uniform(-3, 3), 2) for _ in range(20)]
    return {
        "symbol": symbol,
        "price": price,
        "change": round(price * change_pct / 100, 2),
        "change_pct": change_pct,
        "sparkline": spark,
        "is_live": False,
    }


async def _live_quote(sym: str) -> Optional[dict]:
    """Use Finnhub (preferred) for real-time quote + sparkline."""
    if FINNHUB_KEY:
        try:
            async with httpx.AsyncClient(timeout=15) as c:
                # Quote
                r = await c.get("https://finnhub.io/api/v1/quote",
                                params={"symbol": sym, "token": FINNHUB_KEY})
                r.raise_for_status()
                q = r.json()
                price = float(q.get("c") or 0)
                if not price:
                    raise ValueError("no price")
                change = float(q.get("d") or 0)
                change_pct = float(q.get("dp") or 0)
                prev_close = float(q.get("pc") or price)
                open_p = float(q.get("o") or price)
                high = float(q.get("h") or price)
                low = float(q.get("l") or price)
                # Build a synthetic 20-pt sparkline using today's o/h/l/c around prev close
                spark = [prev_close, open_p, low, (open_p + high) / 2, high, (high + price) / 2,
                         price, (price + low) / 2, low, (low + high) / 2, price, high, low,
                         (price + open_p) / 2, price, high, (high + low) / 2, price, open_p, price]
                return {
                    "symbol": sym,
                    "price": round(price, 2),
                    "change": round(change, 2),
                    "change_pct": round(change_pct, 2),
                    "sparkline": [round(x, 2) for x in spark],
                    "is_live": True,
                }
        except Exception as e:
            logger.warning(f"Finnhub quote fail {sym}: {e}")
    # Fallback to Alpha Vantage
    try:
        gq = await _av_get({"function": "GLOBAL_QUOTE", "symbol": sym})
        q = gq.get("Global Quote") or {}
        price = float(q.get("05. price", 0))
        change = float(q.get("09. change", 0))
        cp = q.get("10. change percent", "0%").replace("%", "")
        change_pct = float(cp)
        if not price:
            return None
        ds = await _av_get({"function": "TIME_SERIES_DAILY", "symbol": sym, "outputsize": "compact"})
        series = ds.get("Time Series (Daily)") or {}
        keys = sorted(series.keys())[-20:]
        spark = [float(series[k]["4. close"]) for k in keys] if keys else [price] * 20
        return {
            "symbol": sym, "price": round(price, 2), "change": round(change, 2),
            "change_pct": round(change_pct, 2), "sparkline": spark, "is_live": True,
        }
    except Exception as e:
        logger.warning(f"AV fallback fail {sym}: {e}")
        return None


async def _live_candles(sym: str, resolution: str = "60", days: int = 1) -> Optional[List[dict]]:
    """Finnhub candles for charting. resolution: 1,5,15,30,60,D,W,M."""
    if not FINNHUB_KEY:
        return None
    try:
        import time as _t
        now = int(_t.time())
        past = now - days * 86400
        async with httpx.AsyncClient(timeout=15) as c:
            r = await c.get("https://finnhub.io/api/v1/stock/candle",
                            params={"symbol": sym, "resolution": resolution,
                                    "from": past, "to": now, "token": FINNHUB_KEY})
            r.raise_for_status()
            j = r.json()
        if j.get("s") != "ok" or not j.get("t"):
            return None
        return [{"t": str(j["t"][i]), "price": float(j["c"][i])} for i in range(len(j["t"]))]
    except Exception as e:
        logger.warning(f"Finnhub candles fail {sym}: {e}")
        return None


@api.get("/stocks/quotes")
async def get_quotes(symbols: str = ",".join(DEFAULT_WATCHLIST)):
    """Returns batch quotes with sparkline. symbols=AAPL,TSLA,..."""
    syms = [s.strip().upper() for s in symbols.split(",") if s.strip()]
    out = []
    invalid = []
    for sym in syms:
        live = await _live_quote(sym)
        if live:
            out.append(live)
        else:
            invalid.append(sym)
    return {"quotes": out, "invalid": invalid}


class ValidateSymbol(BaseModel):
    symbol: str


@api.post("/stocks/validate")
async def validate_symbol(s: ValidateSymbol):
    """Check whether a ticker exists & is tradable via Finnhub."""
    sym = s.symbol.strip().upper()
    if not sym:
        return {"symbol": sym, "valid": False, "reason": "Empty"}
    if not FINNHUB_KEY:
        return {"symbol": sym, "valid": True, "reason": "no key (assume ok)"}
    try:
        async with httpx.AsyncClient(timeout=10) as c:
            r = await c.get("https://finnhub.io/api/v1/quote",
                            params={"symbol": sym, "token": FINNHUB_KEY})
            r.raise_for_status()
            q = r.json()
        price = float(q.get("c") or 0)
        return {"symbol": sym, "valid": price > 0, "price": price,
                "reason": "OK" if price > 0 else "No price returned (likely invalid ticker)"}
    except Exception as e:
        return {"symbol": sym, "valid": False, "reason": str(e)}


@api.get("/stocks/intraday/{symbol}")
async def get_intraday(symbol: str):
    """Hourly view via Finnhub 60min candles."""
    sym = symbol.upper()
    pts = await _live_candles(sym, resolution="60", days=2)
    if pts:
        return {"symbol": sym, "interval": "60min", "points": pts, "is_live": True}
    # AV fallback
    try:
        data = await _av_get({"function": "TIME_SERIES_DAILY", "symbol": sym, "outputsize": "compact"})
        series = data.get("Time Series (Daily)") or {}
        if not series:
            raise ValueError("No data")
        keys = sorted(series.keys())[-30:]
        return {"symbol": sym, "interval": "daily",
                "points": [{"t": k, "price": float(series[k]["4. close"])} for k in keys],
                "is_live": True}
    except Exception as e:
        logger.warning(f"intraday fail {sym}: {e}")
        mock = _mock_quote(sym)
        return {"symbol": sym, "interval": "daily",
                "points": [{"t": f"day-{i}", "price": p} for i, p in enumerate(mock["sparkline"])],
                "is_live": False}


@api.get("/stocks/quarterly/{symbol}")
async def get_quarterly(symbol: str):
    sym = symbol.upper()
    try:
        data = await _av_get({"function": "TIME_SERIES_WEEKLY", "symbol": sym})
        series = data.get("Weekly Time Series") or {}
        if not series:
            raise ValueError("no data")
        keys = sorted(series.keys())[-13:]  # ~quarter (13 weeks)
        points = [{"t": k, "price": float(series[k]["4. close"])} for k in keys]
        return {"symbol": sym, "interval": "weekly", "points": points, "is_live": True}
    except Exception as e:
        logger.warning(f"quarterly fail {sym}: {e}")
        mock = _mock_quote(sym)
        return {
            "symbol": sym,
            "interval": "weekly",
            "points": [{"t": f"week-{i}", "price": p} for i, p in enumerate(mock["sparkline"])],
            "is_live": False,
        }


@api.get("/stocks/news/{symbol}")
async def get_news(symbol: str):
    sym = symbol.upper()
    try:
        data = await _av_get({
            "function": "NEWS_SENTIMENT",
            "tickers": sym,
            "limit": "10",
        })
        feed = data.get("feed", [])[:10]
        items = [{
            "title": f.get("title"),
            "summary": f.get("summary"),
            "source": f.get("source"),
            "url": f.get("url"),
            "time": f.get("time_published"),
            "sentiment": f.get("overall_sentiment_label", "Neutral"),
        } for f in feed]
        if not items:
            raise ValueError("empty")
        return {"symbol": sym, "items": items, "is_live": True}
    except Exception as e:
        logger.warning(f"news fail {sym}: {e}")
        # mock fallback
        return {
            "symbol": sym,
            "items": [
                {"title": f"{sym} extends gains amid sector rotation",
                 "summary": "Analysts highlight strong fundamentals and positive guidance for the quarter ahead.",
                 "source": "Market Watch", "url": "", "time": "20260210T120000", "sentiment": "Bullish"},
                {"title": f"Macro tailwinds support {sym} rally",
                 "summary": "Lower inflation prints and dovish fed commentary lift large-cap tech.",
                 "source": "Reuters", "url": "", "time": "20260210T100000", "sentiment": "Bullish"},
                {"title": f"Options activity on {sym} signals caution",
                 "summary": "Elevated put/call ratio suggests hedging into earnings season.",
                 "source": "Bloomberg", "url": "", "time": "20260210T080000", "sentiment": "Neutral"},
            ],
            "is_live": False,
        }


class RecRequest(BaseModel):
    symbol: str


@api.post("/stocks/recommendation")
async def get_recommendation(req: RecRequest):
    sym = req.symbol.upper()
    # Pull live snapshot
    quotes = await get_quotes(sym)
    q = quotes["quotes"][0]
    news = await get_news(sym)
    headlines = "\n".join(f"- {n['title']} ({n['sentiment']})" for n in news["items"][:5])

    prompt = (
        f"Provide a brief tactical recommendation for {sym}.\n"
        f"Price: ${q['price']}, change today: {q['change_pct']}%.\n"
        f"Recent headlines:\n{headlines}\n\n"
        "Respond in this exact JSON format with no extra text: "
        '{"action": "BUY|HOLD|SELL", "confidence": 0-100, "horizon": "short|medium|long", '
        '"reasoning": "2-3 sentence rationale acknowledging risk"}'
    )

    try:
        chat = LlmChat(
            api_key=EMERGENT_LLM_KEY,
            session_id=f"rec-{sym}-{uuid.uuid4().hex[:6]}",
            system_message="You are JARVIS, a calm tactical market analyst. Always return strict JSON.",
        ).with_model("anthropic", CLAUDE_MODEL)
        reply = await chat.send_message(UserMessage(text=prompt))
        match = re.search(r"\{.*\}", reply, re.S)
        rec = json.loads(match.group(0)) if match else {}
    except Exception as e:
        logger.warning(f"rec ai fail: {e}")
        rec = {"action": "HOLD", "confidence": 60, "horizon": "medium",
               "reasoning": "Insufficient signal; maintain position and monitor headlines."}

    rec["symbol"] = sym
    rec["snapshot"] = q
    return rec


# ==================================================================
# Document Generation (PDF & PPTX)
# ==================================================================
class DocRequest(BaseModel):
    prompt: str
    format: Literal["pdf", "pptx"]
    title: Optional[str] = None


class DocRecord(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    title: str
    format: str
    prompt: str
    created_at: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))
    size_bytes: int


async def _ai_outline(prompt: str, fmt: str) -> dict:
    """Ask Claude for structured plaintext (NOT JSON) — far more reliable."""
    if fmt == "pptx":
        instructions = (
            f"User prompt:\n\"\"\"{prompt}\"\"\"\n\n"
            "Produce a PowerPoint outline of 5-8 slides. "
            "Output STRICT plain text in this exact format (no JSON, no markdown fences):\n\n"
            "TITLE: <document title>\n"
            "SUBTITLE: <one line>\n"
            "\n"
            "SLIDE: <slide 1 title>\n"
            "- <bullet point 1>\n"
            "- <bullet point 2>\n"
            "- <bullet point 3>\n"
            "\n"
            "SLIDE: <slide 2 title>\n"
            "- <bullet>\n"
            "...\n"
            "End with no trailing text. Each slide has 3-5 bullet points."
        )
    else:
        instructions = (
            f"User prompt:\n\"\"\"{prompt}\"\"\"\n\n"
            "Produce a PDF report outline. "
            "Output STRICT plain text in this exact format (no JSON, no markdown fences):\n\n"
            "TITLE: <document title>\n"
            "SUBTITLE: <one line>\n"
            "SUMMARY: <one paragraph executive summary>\n"
            "\n"
            "SECTION: <section 1 heading>\n"
            "PARA: <first paragraph>\n"
            "PARA: <second paragraph>\n"
            "\n"
            "SECTION: <section 2 heading>\n"
            "PARA: <paragraph>\n"
            "...\n"
            "Use 4-6 sections, each 2-4 short paragraphs. No trailing notes."
        )

    last_err = None
    for attempt in range(3):
        try:
            chat = LlmChat(
                api_key=EMERGENT_LLM_KEY,
                session_id=f"doc-{uuid.uuid4().hex[:8]}",
                system_message="You are JARVIS, a content architect. Follow the requested plain-text format exactly. No code fences, no JSON.",
            ).with_model("anthropic", CLAUDE_MODEL)
            reply = await chat.send_message(UserMessage(text=instructions))
            return _parse_marker_outline(reply, fmt)
        except Exception as e:
            last_err = e
            logger.warning(f"outline gen failed attempt {attempt+1}: {e}")
            continue
    raise HTTPException(500, f"Could not generate outline after 3 tries: {last_err}")


def _parse_marker_outline(text: str, fmt: str) -> dict:
    """Parse the TITLE/SLIDE/SECTION/PARA marker format. Robust to quotes & punctuation."""
    lines = [l.strip() for l in text.splitlines()]
    out: dict = {"title": "", "subtitle": "", "summary": "", "slides": [], "sections": []}
    cur_slide = None
    cur_section = None
    for line in lines:
        if not line:
            continue
        up = line.upper()
        if up.startswith("TITLE:"):
            out["title"] = line.split(":", 1)[1].strip()
            cur_slide = None; cur_section = None
        elif up.startswith("SUBTITLE:"):
            out["subtitle"] = line.split(":", 1)[1].strip()
        elif up.startswith("SUMMARY:"):
            out["summary"] = line.split(":", 1)[1].strip()
        elif up.startswith("SLIDE:"):
            cur_slide = {"title": line.split(":", 1)[1].strip(), "bullets": []}
            out["slides"].append(cur_slide)
            cur_section = None
        elif up.startswith("SECTION:"):
            cur_section = {"heading": line.split(":", 1)[1].strip(), "paragraphs": []}
            out["sections"].append(cur_section)
            cur_slide = None
        elif line.startswith("- ") and cur_slide is not None:
            cur_slide["bullets"].append(line[2:].strip())
        elif up.startswith("PARA:") and cur_section is not None:
            cur_section["paragraphs"].append(line.split(":", 1)[1].strip())
        elif cur_section is not None and line:
            # Append continuation to last paragraph
            if cur_section["paragraphs"]:
                cur_section["paragraphs"][-1] += " " + line
            else:
                cur_section["paragraphs"].append(line)
        elif cur_slide is not None and line:
            # Treat plain line as bullet
            cur_slide["bullets"].append(line)
    if not out["title"]:
        out["title"] = "Untitled Document"
    return out


def _build_pdf(outline: dict) -> bytes:
    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=LETTER, leftMargin=0.7*inch, rightMargin=0.7*inch,
                            topMargin=0.7*inch, bottomMargin=0.7*inch)
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("Title", parent=styles["Title"], fontSize=26,
                                 textColor=HexColor("#000000"), spaceAfter=12, alignment=TA_LEFT)
    sub_style = ParagraphStyle("Sub", parent=styles["Italic"], fontSize=13,
                               textColor=HexColor("#555555"), spaceAfter=20)
    h_style = ParagraphStyle("H", parent=styles["Heading2"], fontSize=16,
                             textColor=HexColor("#0B6E4F"), spaceBefore=12, spaceAfter=6)
    body_style = ParagraphStyle("B", parent=styles["BodyText"], fontSize=11, leading=16, spaceAfter=8)

    story = [
        Paragraph(outline.get("title", "Project Inception Report"), title_style),
        Paragraph(outline.get("subtitle", ""), sub_style),
    ]
    if outline.get("summary"):
        story += [Paragraph("Executive Summary", h_style),
                  Paragraph(outline["summary"], body_style)]
    for section in outline.get("sections", []):
        story.append(Paragraph(section.get("heading", ""), h_style))
        for p in section.get("paragraphs", []):
            story.append(Paragraph(p, body_style))
    story.append(Spacer(1, 0.3*inch))
    story.append(Paragraph("<i>Generated by Project Inception &middot; JARVIS</i>",
                           ParagraphStyle("F", parent=styles["BodyText"], fontSize=9,
                                          textColor=HexColor("#888888"), alignment=TA_CENTER)))
    doc.build(story)
    return buf.getvalue()


def _build_pptx(outline: dict) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    BG = RGBColor(0x05, 0x05, 0x05)
    NEON = RGBColor(0xD4, 0xFF, 0x00)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    GRAY = RGBColor(0xA1, 0xA1, 0xAA)

    def add_bg(slide):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = BG

    # Title slide
    s = prs.slides.add_slide(blank)
    add_bg(s)
    tb = s.shapes.add_textbox(Inches(0.6), Inches(2.5), Inches(12), Inches(2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = outline.get("title", "Project Inception")
    r.font.size = Pt(54); r.font.bold = True; r.font.color.rgb = NEON
    p2 = tf.add_paragraph()
    r2 = p2.add_run(); r2.text = outline.get("subtitle", "")
    r2.font.size = Pt(22); r2.font.color.rgb = GRAY

    # Content slides
    for slide_data in outline.get("slides", []):
        s = prs.slides.add_slide(blank); add_bg(s)
        tb = s.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(1))
        tf = tb.text_frame
        r = tf.paragraphs[0].add_run()
        r.text = slide_data.get("title", "")
        r.font.size = Pt(34); r.font.bold = True; r.font.color.rgb = NEON

        body = s.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12), Inches(5.5))
        bf = body.text_frame; bf.word_wrap = True
        first = True
        for bullet in slide_data.get("bullets", []):
            p = bf.paragraphs[0] if first else bf.add_paragraph()
            first = False
            run = p.add_run()
            run.text = f"•  {bullet}"
            run.font.size = Pt(20)
            run.font.color.rgb = WHITE
            p.space_after = Pt(10)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


@api.post("/documents/generate")
async def generate_document(req: DocRequest):
    if not EMERGENT_LLM_KEY:
        raise HTTPException(500, "EMERGENT_LLM_KEY missing")

    outline = await _ai_outline(req.prompt, req.format)
    if req.format == "pdf":
        file_bytes = _build_pdf(outline)
    else:
        file_bytes = _build_pptx(outline)

    title = req.title or outline.get("title") or "Untitled Document"
    record = DocRecord(title=title, format=req.format, prompt=req.prompt,
                       size_bytes=len(file_bytes))
    doc_dict = record.model_dump()
    # Store file as base64 in mongo (small files OK for MVP)
    doc_dict["file_b64"] = base64.b64encode(file_bytes).decode("ascii")
    await db.documents.insert_one(doc_dict)

    return {
        "id": record.id,
        "title": title,
        "format": req.format,
        "size_bytes": record.size_bytes,
        "outline": outline,
        "file_b64": doc_dict["file_b64"],
        "created_at": record.created_at.isoformat(),
    }


@api.get("/documents")
async def list_documents():
    docs = await db.documents.find(
        {}, {"_id": 0, "file_b64": 0}
    ).sort("created_at", -1).to_list(200)
    for d in docs:
        if isinstance(d.get("created_at"), datetime):
            d["created_at"] = d["created_at"].isoformat()
    return {"documents": docs}


@api.get("/documents/{doc_id}")
async def get_document(doc_id: str):
    d = await db.documents.find_one({"id": doc_id}, {"_id": 0})
    if not d:
        raise HTTPException(404, "Not found")
    if isinstance(d.get("created_at"), datetime):
        d["created_at"] = d["created_at"].isoformat()
    return d


@api.delete("/documents/{doc_id}")
async def delete_document(doc_id: str):
    r = await db.documents.delete_one({"id": doc_id})
    return {"deleted": r.deleted_count}


# ==================================================================
# World Crisis Feed — GDELT 2.0 (free, real-time global events)
# Surfaces civil unrest, military action, guerrilla fighting, crises.
# ==================================================================
CRISIS_KEYWORDS = {
    "military": ["airstrike", "missile", "troops", "invasion", "offensive", "military",
                 "army", "naval", "drone strike", "ceasefire", "warplane", "armored",
                 "warship", "submarine", "deployment", "soldiers"],
    "unrest": ["protest", "riot", "demonstration", "unrest", "uprising", "march",
               "rally", "civil disobedience", "general strike", "clash with police"],
    "guerrilla": ["insurgent", "guerrilla", "guerilla", "militants", "militia",
                  "rebels", "armed group", "ambush", "terror", "terrorist", "extremist",
                  "isis", "isil", "al-qaeda", "boko haram", "taliban", "houthis",
                  "hezbollah", "hamas", "wagner", "junta"],
    "crisis": ["famine", "earthquake", "hurricane", "typhoon", "refugee", "evacuation",
               "humanitarian", "state of emergency", "wildfire", "flood disaster",
               "outbreak", "epidemic", "sanctions", "coup", "crisis"],
}
HOTSPOT_TERMS = [
    "Gaza", "Israel", "Lebanon", "Syria", "Yemen", "Iran", "Iraq",
    "Sudan", "Somalia", "Mali", "Burkina Faso", "Niger", "Ethiopia",
    "Ukraine", "Russia", "Belarus", "Moldova",
    "Afghanistan", "Pakistan", "Kashmir",
    "Myanmar", "Taiwan", "North Korea", "South China Sea",
    "Haiti", "Venezuela", "Colombia", "Mexico cartel",
    "Nigeria", "DR Congo", "Central African Republic", "Mozambique",
    "Libya", "West Bank", "Sinai", "Tigray", "Donbas", "Kursk",
]


def _classify(text: str) -> tuple[str, list[str]]:
    """Return (primary_category, matched_hotspots). primary='other' if none match."""
    t = (text or "").lower()
    matched = [h for h in HOTSPOT_TERMS if h.lower() in t]
    for cat in ("military", "guerrilla", "unrest", "crisis"):
        if any(k in t for k in CRISIS_KEYWORDS[cat]):
            return cat.upper(), matched
    return "OTHER", matched


_crisis_cache: dict = {}


@api.get("/world/crisis")
async def world_crisis(category: str = "all", limit: int = 30):
    cat = (category or "all").lower()
    key = f"{cat}:{limit}"
    now_ts = datetime.now(timezone.utc).timestamp()
    cached = _crisis_cache.get(key)
    if cached and (now_ts - cached["t"]) < 90:
        return cached["data"]

    raw_items: list[dict] = []
    # 1) Finnhub general news (real-time breaking world news)
    if FINNHUB_KEY:
        try:
            async with httpx.AsyncClient(timeout=12) as c:
                r = await c.get("https://finnhub.io/api/v1/news",
                                params={"category": "general", "token": FINNHUB_KEY})
                r.raise_for_status()
                for a in (r.json() or []):
                    raw_items.append({
                        "title": a.get("headline") or "",
                        "summary": a.get("summary") or "",
                        "url": a.get("url") or "",
                        "source": a.get("source") or "",
                        "image": a.get("image") or "",
                        "time": datetime.fromtimestamp(a.get("datetime", 0),
                                                      tz=timezone.utc).isoformat() if a.get("datetime") else "",
                    })
        except Exception as e:
            logger.warning(f"Finnhub general news fail: {e}")

    # Classify + filter
    classified = []
    for item in raw_items:
        blob = f"{item['title']} {item.get('summary', '')}"
        pri, hotspots = _classify(blob)
        item["category"] = pri
        item["hotspots"] = hotspots
        classified.append(item)

    if cat == "all":
        # Crisis-only across military/guerrilla/unrest/crisis
        filtered = [i for i in classified if i["category"] != "OTHER"]
    else:
        filtered = [i for i in classified if i["category"] == cat.upper()]

    # If thin, supplement with hotspot-matched items
    if len(filtered) < 5:
        for i in classified:
            if i not in filtered and i["hotspots"]:
                filtered.append(i)
                if len(filtered) >= limit:
                    break

    # Always also include hotspot-tagged stories even if no crisis keyword
    out = {"category": cat, "count": len(filtered[:limit]), "items": filtered[:limit]}
    _crisis_cache[key] = {"t": now_ts, "data": out}
    return out


# Hotspot timezone directory
HOTSPOT_ZONES = [
    {"name": "Gaza City",       "zone": "Asia/Gaza",       "region": "Middle East"},
    {"name": "Tel Aviv",        "zone": "Asia/Jerusalem",  "region": "Middle East"},
    {"name": "Beirut",          "zone": "Asia/Beirut",     "region": "Middle East"},
    {"name": "Damascus",        "zone": "Asia/Damascus",   "region": "Middle East"},
    {"name": "Sana'a",          "zone": "Asia/Aden",       "region": "Middle East"},
    {"name": "Tehran",          "zone": "Asia/Tehran",     "region": "Middle East"},
    {"name": "Baghdad",         "zone": "Asia/Baghdad",    "region": "Middle East"},
    {"name": "Kabul",           "zone": "Asia/Kabul",      "region": "South Asia"},
    {"name": "Islamabad",       "zone": "Asia/Karachi",    "region": "South Asia"},
    {"name": "Kyiv",            "zone": "Europe/Kyiv",     "region": "Eastern Europe"},
    {"name": "Moscow",          "zone": "Europe/Moscow",   "region": "Eastern Europe"},
    {"name": "Minsk",           "zone": "Europe/Minsk",    "region": "Eastern Europe"},
    {"name": "Khartoum",        "zone": "Africa/Khartoum", "region": "East Africa"},
    {"name": "Mogadishu",       "zone": "Africa/Mogadishu","region": "East Africa"},
    {"name": "Addis Ababa",     "zone": "Africa/Addis_Ababa","region": "East Africa"},
    {"name": "Tripoli",         "zone": "Africa/Tripoli",  "region": "North Africa"},
    {"name": "Bamako",          "zone": "Africa/Bamako",   "region": "West Africa"},
    {"name": "Ouagadougou",     "zone": "Africa/Ouagadougou","region": "West Africa"},
    {"name": "Abuja",           "zone": "Africa/Lagos",    "region": "West Africa"},
    {"name": "Kinshasa",        "zone": "Africa/Kinshasa", "region": "Central Africa"},
    {"name": "Naypyidaw",       "zone": "Asia/Yangon",     "region": "SE Asia"},
    {"name": "Pyongyang",       "zone": "Asia/Pyongyang",  "region": "East Asia"},
    {"name": "Taipei",          "zone": "Asia/Taipei",     "region": "East Asia"},
    {"name": "Port-au-Prince",  "zone": "America/Port-au-Prince","region": "Caribbean"},
    {"name": "Caracas",         "zone": "America/Caracas", "region": "South America"},
    {"name": "Bogota",          "zone": "America/Bogota",  "region": "South America"},
]


@api.get("/world/hotspots")
async def world_hotspots():
    """Returns hotspot cities with current local time. Frontend renders relative clock."""
    from zoneinfo import ZoneInfo
    now_utc = datetime.now(timezone.utc)
    out = []
    for h in HOTSPOT_ZONES:
        try:
            local = now_utc.astimezone(ZoneInfo(h["zone"]))
            offset_minutes = int(local.utcoffset().total_seconds() / 60) if local.utcoffset() else 0
            offset_str = ("+" if offset_minutes >= 0 else "-") + \
                f"{abs(offset_minutes)//60:02d}:{abs(offset_minutes)%60:02d}"
            out.append({
                **h,
                "local_time": local.strftime("%H:%M"),
                "local_date": local.strftime("%Y-%m-%d"),
                "offset": offset_str,
                "offset_minutes": offset_minutes,
                "weekday": local.strftime("%a").upper(),
            })
        except Exception as e:
            logger.warning(f"tz fail {h['zone']}: {e}")
    return {"server_utc": now_utc.isoformat(), "hotspots": out}


# ==================================================================
# TTS / STT  (OpenAI via Emergent gateway)
# ==================================================================
class TTSRequest(BaseModel):
    text: str
    voice: str = "nova"  # alloy, echo, fable, onyx, nova, shimmer


@api.post("/tts")
async def tts(req: TTSRequest):
    """Returns base64 mp3 for narration. Used by accessibility/voice mode."""
    try:
        resp = await oai_client.audio.speech.create(
            model="tts-1",
            voice=req.voice,
            input=req.text[:4000],
        )
        audio_bytes = await resp.aread() if hasattr(resp, "aread") else resp.content
        return {"audio_b64": base64.b64encode(audio_bytes).decode("ascii"), "mime": "audio/mpeg"}
    except Exception as e:
        logger.exception("tts fail")
        raise HTTPException(500, f"TTS error: {e}")


@api.post("/stt")
async def stt(file: UploadFile = File(...)):
    """Transcribes uploaded audio via Whisper."""
    try:
        content = await file.read()
        # OpenAI SDK accepts file-like with name
        bio = io.BytesIO(content)
        bio.name = file.filename or "audio.m4a"
        result = await oai_client.audio.transcriptions.create(
            model="whisper-1",
            file=bio,
        )
        text = result.text if hasattr(result, "text") else str(result)
        return {"text": text}
    except Exception as e:
        logger.exception("stt fail")
        raise HTTPException(500, f"STT error: {e}")


# ==================================================================
# Settings (watchlist persistence per device/user id)
# ==================================================================
class Watchlist(BaseModel):
    user_id: str
    symbols: List[str]


@api.get("/watchlist/{user_id}")
async def get_watchlist(user_id: str):
    d = await db.watchlists.find_one({"user_id": user_id}, {"_id": 0})
    if not d:
        return {"user_id": user_id, "symbols": DEFAULT_WATCHLIST}
    return d


@api.put("/watchlist")
async def set_watchlist(w: Watchlist):
    await db.watchlists.update_one(
        {"user_id": w.user_id},
        {"$set": w.model_dump()},
        upsert=True,
    )
    return w.model_dump()


# ==================================================================
# DAGRCMD — Encrypted Comms (E2E)
# Server NEVER sees plaintext. Each message is encrypted client-side using
# tweetnacl (X25519 + XSalsa20-Poly1305) once per recipient.
# Server stores only ciphertext maps + nonces and relays via WebSocket.
# ==================================================================
import hashlib
import secrets

def _hash_code(callsign: str, code: str) -> str:
    """Hash auth code with callsign as salt (NEVER store raw)."""
    return hashlib.sha256(f"{callsign.lower()}::{code}".encode()).hexdigest()


class OfficerRegister(BaseModel):
    callsign: str
    auth_code: str  # 6+ char password / pin
    public_key: str  # base64 X25519 public key
    rank: Optional[str] = None
    unit: Optional[str] = None


class OfficerLogin(BaseModel):
    callsign: str
    auth_code: str


@api.post("/dagrcmd/officers/register")
async def register_officer(o: OfficerRegister):
    callsign = o.callsign.strip().upper()
    if not callsign or not o.auth_code or not o.public_key:
        raise HTTPException(400, "callsign, auth_code, public_key required")
    existing = await db.dagr_officers.find_one({"callsign": callsign})
    if existing:
        # If same auth_code, allow public_key rotation (re-install)
        if existing.get("auth_hash") != _hash_code(callsign, o.auth_code):
            raise HTTPException(409, "Callsign already exists")
        await db.dagr_officers.update_one(
            {"callsign": callsign},
            {"$set": {"public_key": o.public_key,
                      "rank": o.rank or existing.get("rank"),
                      "unit": o.unit or existing.get("unit")}}
        )
        return {"callsign": callsign, "rotated": True}
    doc = {
        "callsign": callsign,
        "auth_hash": _hash_code(callsign, o.auth_code),
        "public_key": o.public_key,
        "rank": o.rank or "OPERATOR",
        "unit": o.unit or "UNASSIGNED",
        "created_at": datetime.now(timezone.utc).isoformat(),
        "last_seen": datetime.now(timezone.utc).isoformat(),
    }
    await db.dagr_officers.insert_one(doc)
    return {"callsign": callsign, "rotated": False}


@api.post("/dagrcmd/officers/login")
async def login_officer(o: OfficerLogin):
    callsign = o.callsign.strip().upper()
    doc = await db.dagr_officers.find_one({"callsign": callsign}, {"_id": 0})
    if not doc or doc.get("auth_hash") != _hash_code(callsign, o.auth_code):
        raise HTTPException(401, "Invalid callsign or auth code")
    await db.dagr_officers.update_one(
        {"callsign": callsign},
        {"$set": {"last_seen": datetime.now(timezone.utc).isoformat()}}
    )
    doc.pop("auth_hash", None)
    return doc


@api.get("/dagrcmd/officers")
async def list_officers(callsigns: Optional[str] = None):
    """Returns directory of officers (callsign + public_key). For key exchange."""
    q = {}
    if callsigns:
        wanted = [c.strip().upper() for c in callsigns.split(",") if c.strip()]
        q = {"callsign": {"$in": wanted}}
    docs = await db.dagr_officers.find(q, {"_id": 0, "auth_hash": 0}).to_list(500)
    return {"officers": docs}


class ChannelCreate(BaseModel):
    name: str
    owner: str  # callsign
    auth_code: str
    members: List[str] = []


@api.post("/dagrcmd/channels")
async def create_channel(c: ChannelCreate):
    owner = c.owner.strip().upper()
    owner_doc = await db.dagr_officers.find_one({"callsign": owner})
    if not owner_doc or owner_doc.get("auth_hash") != _hash_code(owner, c.auth_code):
        raise HTTPException(401, "Auth failed")
    members = sorted({owner, *(m.strip().upper() for m in c.members)})
    channel = {
        "id": str(uuid.uuid4()),
        "name": c.name.strip()[:64] or "CHANNEL",
        "owner": owner,
        "members": members,
        "join_code": secrets.token_hex(3).upper(),  # 6-char invite code
        "created_at": datetime.now(timezone.utc).isoformat(),
    }
    await db.dagr_channels.insert_one(channel)
    channel.pop("_id", None)
    return channel


@api.get("/dagrcmd/channels/{callsign}")
async def list_channels(callsign: str):
    cs = callsign.strip().upper()
    docs = await db.dagr_channels.find({"members": cs}, {"_id": 0}).sort("created_at", -1).to_list(200)
    return {"channels": docs}


class ChannelJoin(BaseModel):
    callsign: str
    auth_code: str
    join_code: str


@api.post("/dagrcmd/channels/join")
async def join_channel(j: ChannelJoin):
    cs = j.callsign.strip().upper()
    code = j.join_code.strip().upper()
    officer = await db.dagr_officers.find_one({"callsign": cs})
    if not officer or officer.get("auth_hash") != _hash_code(cs, j.auth_code):
        raise HTTPException(401, "Auth failed")
    chan = await db.dagr_channels.find_one({"join_code": code}, {"_id": 0})
    if not chan:
        raise HTTPException(404, "Invalid join code")
    if cs not in chan["members"]:
        await db.dagr_channels.update_one(
            {"id": chan["id"]}, {"$addToSet": {"members": cs}}
        )
        chan["members"] = sorted(set(chan["members"] + [cs]))
    return chan


class EncryptedMessage(BaseModel):
    channel_id: str
    sender: str
    sender_pubkey: str
    kind: Literal["text", "audio", "location"] = "text"
    # ciphertexts: {recipient_callsign: {"ct": "...", "nonce": "..."}}
    ciphertexts: dict
    meta: Optional[dict] = None  # optional non-sensitive metadata (e.g., audio_ms)


@api.post("/dagrcmd/messages")
async def send_message(m: EncryptedMessage):
    chan = await db.dagr_channels.find_one({"id": m.channel_id}, {"_id": 0})
    if not chan:
        raise HTTPException(404, "Channel not found")
    if m.sender not in chan["members"]:
        raise HTTPException(403, "Sender not in channel")
    msg = {
        "id": str(uuid.uuid4()),
        "channel_id": m.channel_id,
        "sender": m.sender,
        "sender_pubkey": m.sender_pubkey,
        "kind": m.kind,
        "ciphertexts": m.ciphertexts,
        "meta": m.meta or {},
        "timestamp": datetime.now(timezone.utc).isoformat(),
    }
    await db.dagr_messages.insert_one(msg)
    msg.pop("_id", None)
    # Broadcast to connected members
    await _dagr_broadcast(chan["members"], {"type": "message", "data": msg})
    return msg


@api.get("/dagrcmd/messages/{channel_id}")
async def list_messages(channel_id: str, callsign: str, limit: int = 100):
    cs = callsign.strip().upper()
    chan = await db.dagr_channels.find_one({"id": channel_id})
    if not chan or cs not in chan["members"]:
        raise HTTPException(403, "Not a member")
    docs = await db.dagr_messages.find(
        {"channel_id": channel_id}, {"_id": 0}
    ).sort("timestamp", -1).to_list(limit)
    docs.reverse()
    # Only return ciphertext for this recipient (privacy) — keep size small
    out = []
    for d in docs:
        ct = d["ciphertexts"].get(cs)
        if not ct and d["sender"] != cs:
            continue
        out.append({**d, "ciphertext_for_me": ct})
    return {"messages": out}


# ---------- WebSocket relay (real-time) ----------
class _Hub:
    def __init__(self):
        self.conns: dict = {}  # callsign -> WebSocket

    async def connect(self, callsign: str, ws: WebSocket):
        await ws.accept()
        # Drop prior connection if any
        old = self.conns.get(callsign)
        if old:
            try: await old.close()
            except Exception: pass
        self.conns[callsign] = ws

    def disconnect(self, callsign: str, ws: WebSocket):
        if self.conns.get(callsign) is ws:
            self.conns.pop(callsign, None)

    async def send(self, callsign: str, payload: dict):
        ws = self.conns.get(callsign)
        if ws is None:
            return False
        try:
            await ws.send_json(payload)
            return True
        except Exception:
            self.conns.pop(callsign, None)
            return False


hub = _Hub()


async def _dagr_broadcast(callsigns: list, payload: dict):
    for cs in callsigns:
        await hub.send(cs, payload)


@app.websocket("/api/ws/dagrcmd/{callsign}")
async def ws_dagrcmd(websocket: WebSocket, callsign: str, auth_code: str = ""):
    cs = callsign.strip().upper()
    officer = await db.dagr_officers.find_one({"callsign": cs})
    if not officer or officer.get("auth_hash") != _hash_code(cs, auth_code):
        await websocket.close(code=4401)
        return
    await hub.connect(cs, websocket)
    await websocket.send_json({"type": "ready", "callsign": cs,
                               "server_time": datetime.now(timezone.utc).isoformat()})
    try:
        while True:
            data = await websocket.receive_json()
            # Heartbeat / presence pings only — actual encrypted messages go via REST
            if data.get("type") == "ping":
                await websocket.send_json({"type": "pong",
                                          "t": datetime.now(timezone.utc).isoformat()})
            elif data.get("type") == "presence":
                # Broadcast presence to channel members
                channels = await db.dagr_channels.find({"members": cs}).to_list(50)
                everyone = sorted({m for c in channels for m in c["members"]} - {cs})
                await _dagr_broadcast(everyone, {"type": "presence",
                                                 "callsign": cs,
                                                 "status": data.get("status", "online")})
    except WebSocketDisconnect:
        pass
    except Exception as e:
        logger.warning(f"ws err {cs}: {e}")
    finally:
        hub.disconnect(cs, websocket)


# ==================================================================
# Health
# ==================================================================
@api.get("/")
async def root():
    return {"app": "Project Inception", "status": "ok",
            "alpha_vantage_configured": bool(ALPHA_VANTAGE_KEY),
            "llm_configured": bool(EMERGENT_LLM_KEY)}


# Serve Ionicons font directly so the frontend can bypass Metro asset bundling
# (which sometimes returns an empty file under tunnel mode).
IONICONS_TTF = Path("/app/frontend/node_modules/@expo/vector-icons/build/vendor/react-native-vector-icons/Fonts/Ionicons.ttf")


@api.get("/fonts/ionicons.ttf")
async def fonts_ionicons():
    if IONICONS_TTF.exists():
        return FileResponse(str(IONICONS_TTF), media_type="font/ttf")
    raise HTTPException(404, "Font not found")


@app.get("/api/join/{code}", response_class=Response)
async def join_landing(code: str):
    """User-facing join landing page. Shows install instructions + the invite code."""
    code_safe = (code or "").upper().replace("<", "").replace(">", "")[:12]
    chan = await db.dagr_channels.find_one({"join_code": code_safe}, {"_id": 0})
    chan_name = (chan or {}).get("name", "—")
    chan_owner = (chan or {}).get("owner", "—")
    members = len((chan or {}).get("members", [])) if chan else 0
    app_url = "https://inception-ai-gen.preview.emergentagent.com"
    html = f"""<!DOCTYPE html>
<html lang="en"><head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width,initial-scale=1">
<title>DAGRCMD Invite · {code_safe}</title>
<style>
  body {{ margin:0; padding:0; background:#0a0000; color:#FFB6B6;
         font-family:-apple-system,system-ui,Segoe UI,Roboto,sans-serif; }}
  .wrap {{ max-width:520px; margin:0 auto; padding:32px 20px; }}
  h1 {{ color:#FF1A1A; letter-spacing:6px; font-weight:700; font-size:28px;
        margin:0 0 4px; text-align:center; }}
  .sub {{ color:#7A2424; letter-spacing:2px; font-size:11px;
          text-align:center; margin-bottom:24px; }}
  .code {{ background:#140404; border:1px solid #FF1A1A; border-radius:8px;
           padding:20px; text-align:center; margin:16px 0; }}
  .code .label {{ color:#FF1A1A; letter-spacing:3px; font-size:11px; }}
  .code .value {{ color:#FF4444; font-size:42px; letter-spacing:6px;
                  font-weight:700; font-family:monospace; margin-top:8px; }}
  .card {{ background:#140404; border:1px solid #2E0A0A; border-radius:8px;
           padding:18px; margin:12px 0; }}
  .card h3 {{ color:#FF1A1A; letter-spacing:2px; font-size:12px;
              margin:0 0 10px; }}
  .row {{ display:flex; gap:14px; margin:8px 0; align-items:flex-start; }}
  .num {{ flex-shrink:0; width:30px; height:30px; border-radius:15px;
          background:#FF1A1A; color:#0a0000; text-align:center; line-height:30px;
          font-weight:700; font-size:14px; }}
  .row p {{ margin:4px 0; line-height:1.4; font-size:14px; color:#FFB6B6; }}
  .row p b {{ color:#fff; }}
  a {{ color:#FF4444; text-decoration:none; word-break:break-all; }}
  .btn {{ display:block; background:#FF1A1A; color:#0a0000; text-align:center;
          padding:14px; border-radius:8px; font-weight:700; letter-spacing:2px;
          margin:16px 0; font-size:14px; }}
  .small {{ color:#7A2424; font-size:11px; text-align:center; margin-top:24px;
            letter-spacing:1px; }}
</style></head>
<body>
<div class="wrap">
  <h1>[ DAGRCMD ]</h1>
  <div class="sub">CLASSIFIED SECURE TERMINAL</div>

  <div class="card">
    <h3>YOU HAVE BEEN INVITED</h3>
    <p style="margin:6px 0;color:#FFB6B6;font-size:14px">
      <b style="color:#fff">{chan_name}</b><br>
      Owner: {chan_owner} · {members} operators
    </p>
  </div>

  <div class="code">
    <div class="label">INVITE CODE</div>
    <div class="value">{code_safe}</div>
  </div>

  <div class="card">
    <h3>HOW TO JOIN</h3>
    <div class="row"><div class="num">1</div><p>Install <b>Expo Go</b> on your phone:<br>
      <a href="https://apps.apple.com/app/expo-go/id982107779">iOS App Store</a><br>
      <a href="https://play.google.com/store/apps/details?id=host.exp.exponent">Google Play</a></p></div>
    <div class="row"><div class="num">2</div><p>Open <b>Expo Go</b> and scan the QR or enter:<br><a href="{app_url}">{app_url}</a></p></div>
    <div class="row"><div class="num">3</div><p>Tap the red <b>DAGRCMD</b> chip on the home screen, then <b>ENLIST</b> (or AUTHENTICATE if you already have a callsign).</p></div>
    <div class="row"><div class="num">4</div><p>Inside COMMS, tap <b>JOIN BY CODE</b> and paste:<br><b style="color:#FF4444;font-family:monospace;font-size:18px;letter-spacing:3px">{code_safe}</b></p></div>
  </div>

  <a class="btn" href="{app_url}">OPEN DAGRCMD</a>
  <p class="small">END-TO-END ENCRYPTED · X25519 · NACL BOX</p>
</div>
</body></html>"""
    return Response(content=html, media_type="text/html")


app.include_router(api)

app.add_middleware(
    CORSMiddleware,
    allow_credentials=True,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)


@app.on_event("shutdown")
async def _shutdown():
    mongo_client.close()
